# -*- coding: utf-8 -*-
"""WebLogic/Java 이관 검증용 골든 파일 캡처 — 이관 설계 §8의 "단계 0".

현재 Python 시스템의 **결정적** API 응답과 산출물을 떠 둔다. Java 재작성이 끝나면
같은 시나리오를 돌려 이 파일들과 비교한다("같은 입력 → 같은 출력"). 프런트가 안
바뀌므로 이 비교가 통과하면 화면도 통과한다는 것이 이관 설계의 전제다.

실행:  py -3.14 -m golden.capture
  - 임시 디렉터리에 새 DB를 시드하고(운영 data/는 건드리지 않는다) TestClient로
    엔드포인트를 순서대로 때린 뒤, 정규화(아래)를 거쳐 golden/api/*.json에 쓴다.
  - 재실행하면 전부 다시 떠서 덮어쓴다 — git diff가 곧 "무엇이 달라졌나"다.

정규화 규칙(README.md에도 있다 — Java 쪽 비교 하네스도 같은 규칙을 써야 한다):
  - ISO 타임스탬프(초 단위 이상 정밀도)는 값이 실행 시각에 좌우되므로 "<TS>"로 치환.
  - 임시 작업 디렉터리 절대경로는 "<WORK>"로, 리포 루트는 "<REPO>"로 치환.
  - 그 외(정수 id 포함)는 치환하지 않는다 — 시드·시나리오 순서가 고정이라 결정적이다.

캡처하지 않는 것(비결정적 — README에 근거와 함께 기록):
  - LLM 경로 전부: 대화(chat)·업로드 즉시검사·워크플로 실행(run)·작업 메시지.
  - 검색의 벡터 모드: 임베딩 모델 부동소수가 환경 종속 — FTS 모드만 캡처한다.
  - PPTX 바이너리: zip 내부 타임스탬프가 비결정적 — 대신 슬라이드 텍스트를 뜬다.
"""

from __future__ import annotations

import json
import re
import shutil
import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
GOLDEN = REPO / "golden"

sys.path.insert(0, str(REPO))

from fastapi.testclient import TestClient  # noqa: E402

from agent.nodes.pptx_builder import build_pptx  # noqa: E402
from agent.retrieval.indexer import build_index  # noqa: E402
from agent.rfp_text import extract_pdf_text  # noqa: E402
from backend.db import get_connection, init_db  # noqa: E402
from backend.main import create_app  # noqa: E402
from backend.seed import seed_giganlist_districts  # noqa: E402

# 배점표 픽스처 — rfp-locate 스킬의 수원시 정답지를 그대로 쓴다(사람이 검증한 실물).
SCORING_FIXTURE = REPO / ".claude" / "skills" / "rfp-locate" / "references" / "scoring_schema.json"
SUWON_PDF = REPO / "corpus" / "rfp" / "수원시 금고 지정 계획 공고문.pdf"

TS_RE = re.compile(
    r"\d{4}-\d{2}-\d{2}[T ]\d{2}:\d{2}:\d{2}(?:\.\d+)?(?:Z|[+-]\d{2}:?\d{2})?"
)
# 백엔드가 secrets.token_hex(4)로 만드는 랜덤 식별자들(bc-/task-/ntf-/msg-/chat-/new-).
# 실행마다 달라지므로 접두사만 남기고 치환한다 — Java 쪽도 같은 규칙으로 정규화한다.
RANDOM_ID_RE = re.compile(r"\b(bc|task|ntf|msg|chat|new)-[0-9a-f]{8}\b")


def _normalize(obj, work: str):
    if isinstance(obj, dict):
        return {k: _normalize(v, work) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_normalize(v, work) for v in obj]
    if isinstance(obj, str):
        s = TS_RE.sub("<TS>", obj)
        s = RANDOM_ID_RE.sub(lambda m: m.group(1) + "-<ID>", s)
        s = s.replace(work, "<WORK>").replace(work.replace("\\", "/"), "<WORK>")
        s = s.replace(str(REPO), "<REPO>").replace(str(REPO).replace("\\", "/"), "<REPO>")
        return s
    return obj


def main() -> None:
    work = tempfile.mkdtemp(prefix="golden_")
    workp = Path(work)
    api_dir = GOLDEN / "api"
    art_dir = GOLDEN / "artifacts"
    for d in (api_dir, art_dir):
        if d.exists():
            shutil.rmtree(d)
        d.mkdir(parents=True)

    # ── 준비: 새 DB 시드 + 산출물 픽스처 + FTS 인덱스 ─────────────────────
    db_path = workp / "registry.db"
    conn = init_db(str(db_path))
    seeded = seed_giganlist_districts(conn, REPO / "corpus" / "institutions")
    conn.close()

    out_root = workp / "report_new"
    scoring = json.loads(SCORING_FIXTURE.read_text(encoding="utf-8"))
    # coverage_map v2 픽스처 — 병합 응답(coverage-map API)의 골든을 위해.
    nowon_dir = out_root / "노원구"
    nowon_dir.mkdir(parents=True)
    (nowon_dir / "rfp_scoring.json").write_text(
        json.dumps({"institution": "노원구", "rfp_title": scoring["rfp_title"],
                    "total_score": scoring["total_score"], "criteria": scoring["criteria"]},
                   ensure_ascii=False, indent=2), encoding="utf-8")
    (nowon_dir / "coverage_map.json").write_text(json.dumps({
        "version": 2,
        "items": {scoring["criteria"][0]["item"]:
                  {"team": "전산", "covered": True, "gap_note": None}},
        "teams": {"전산": {"pii_count": 3}},
    }, ensure_ascii=False), encoding="utf-8")

    index_db = workp / "corpus_index.db"
    stats = build_index(corpus_root=REPO / "corpus", db_path=index_db, embed=False)

    app = create_app(
        str(db_path),
        output_root=str(out_root),
        index_db_path=str(index_db),
        corpus_root=str(REPO / "corpus"),
        inbox_root=str(workp / "inbox"),
        rfp_root=str(REPO / "corpus" / "rfp"),
        batches_root=str(workp / "batches"),
        graph_db_path=str(workp / "graph.db"),
        archive_root=str(workp / "archive"),
    )
    client = TestClient(app)

    captured: list[str] = []

    def snap(name: str, method: str, url: str, *, body=None, headers=None, expect=None):
        fn = getattr(client, method)
        kwargs = {}
        if body is not None:
            kwargs["json"] = body
        if headers:
            kwargs["headers"] = headers
        resp = fn(url, **kwargs)
        if expect is not None and resp.status_code != expect:
            raise SystemExit(f"[중단] {method.upper()} {url} → {resp.status_code} (기대 {expect}): {resp.text[:300]}")
        try:
            payload = resp.json()
        except ValueError:
            payload = {"_text": resp.text}
        # 요청 쪽 url에도 랜덤 id가 들어간다(/bidcases/bc-xxxx) — 문서 전체를 정규화한다.
        doc = _normalize({
            "request": {"method": method.upper(), "url": url,
                        **({"body": body} if body is not None else {}),
                        **({"headers": headers} if headers else {})},
            "status": resp.status_code,
            "body": payload,
        }, work)
        path = api_dir / f"{len(captured):02d}_{name}.json"
        path.write_text(json.dumps(doc, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        captured.append(name)
        return payload

    from urllib.parse import quote as _q

    # ── 1. 시드 직후의 읽기 API ───────────────────────────────────────────
    insts = snap("institutions_list", "get", "/institutions", expect=200)
    first_id = insts[0]["institution_id"]
    snap("institution_detail", "get", f"/institutions/{first_id}", expect=200)
    snap("institution_404", "get", "/institutions/ghost", expect=404)
    snap("accounts", "get", "/accounts", expect=200)
    snap("menus_default", "get", "/menus", expect=200)
    snap("menus_team_lead", "get", f"/menus?team={_q('전산팀')}&position={_q('팀장')}", expect=200)
    snap("document_read", "get",
         "/documents?path=" + _q("corpus/institutions/dobong/spec/00_인덱스.txt"), expect=200)
    snap("consistency_initial", "get", "/consistency", expect=200)
    snap("coverage_map_nowon", "get", "/institutions/nowon/coverage-map", expect=200)
    snap("artifacts_nowon", "get", "/institutions/nowon/artifacts", expect=200)

    # ── 2. 결정적 쓰기 시나리오: 결재 3단 → 작업 3건 승인 → 최종 확정 ─────
    # 시퀀스는 backend/tests/test_api_bidcases.py의 E2E와 동일하다(같은 입력이어야
    # Java 쪽도 같은 시나리오를 돌릴 수 있다). LLM이 걸리는 작업 대화(messages)는
    # 건너뛴다 — submit/approve는 대화 없이 동작한다.
    bc = snap("bidcase_create", "post", "/bidcases",
              body={"institution_id": first_id, "title": "골든 캡처용 입찰 건",
                    "note": "이관 검증 시나리오 — 값 임의 변경 금지"}, expect=200)
    bid_id = bc["bid_case_id"]
    for tier, by in [(1, "alice"), (2, "bob"), (3, "carol")]:
        snap(f"participation_tier{tier}", "post",
             f"/bidcases/{bid_id}/participation-decisions",
             body={"tier": tier, "role": "영업팀", "by": by, "choice": "참여"}, expect=200)
    detail = snap("bidcase_detail_confirmed", "get", f"/bidcases/{bid_id}", expect=200)
    for n, task in enumerate(detail["tasks"]):
        tid = task["task_id"]
        # 임시저장이 assignee를 선점한다(claim_assignee_if_unset) — 제출은 assignee만
        # 할 수 있다. 테스트가 쓰는 messages POST는 LLM 경로라 여기서는 draft로 선점한다.
        snap(f"task{n}_draft_claim", "patch", f"/tasks/{tid}/draft",
             body={"content": "골든 캡처 고정 메모"}, headers={"X-User-Id": "dave"}, expect=200)
        snap(f"task{n}_submit", "post", f"/tasks/{tid}/submit",
             headers={"X-User-Id": "dave"}, expect=200)
        snap(f"task{n}_approve", "post", f"/tasks/{tid}/approve",
             body={"approved": True}, headers={"X-User-Id": "boss"}, expect=200)
    snap("bidcase_finalize", "post", f"/bidcases/{bid_id}/finalize",
         body={"approved": True}, headers={"X-User-Id": "golden-capture"}, expect=200)
    snap("bidcases_assignee_view", "get", f"/bidcases?team={_q('영업')}&assignee=dave")
    snap("tasks_team_view", "get", f"/tasks?team={_q('전산')}")
    snap("notifications_sales", "get", f"/notifications?recipient={_q('영업팀')}")
    snap("consistency_after", "get", "/consistency", expect=200)
    snap("timeline_after", "get", f"/institutions/{first_id}/timeline", expect=200)
    snap("workflow_status", "get", f"/institutions/{first_id}/status")

    # ── 3. 검색 (FTS 모드만 — 벡터는 임베딩 환경 종속이라 제외) ───────────
    # '금고'(2글자)가 0건인 것은 버그가 아니라 trigram 색인의 원리적 한계다(3글자
    # 미만은 트라이그램을 만들 수 없다). 이 골든이 그 동작을 계약으로 못 박는다 —
    # Java(Oracle Text)가 여기서 결과를 내면 "의도된 개선"으로 문서화할 것(§6-A).
    for n, q in enumerate(["금고", "소상공인 지원", "청년 창업"]):
        snap(f"search_{n}_{q.replace(' ', '_')}", "get", f"/search?q={_q(q)}&limit=5", expect=200)

    # ── 4. 산출물 골든 ───────────────────────────────────────────────────
    # 4-a. PDF 텍스트 추출 — Java(PDFBox)와 비교할 핵심 지점.
    extracted = extract_pdf_text(str(SUWON_PDF))
    (art_dir / "suwon_rfp_text.txt").write_text(extracted["full_text"], encoding="utf-8")
    (art_dir / "suwon_rfp_text.meta.json").write_text(json.dumps({
        "source_pdf": SUWON_PDF.name,
        "pages": len(extracted["pages"]),
        "full_text_chars": len(extracted["full_text"]),
        "avg_chars_per_page": extracted["avg_chars_per_page"],
        "is_abnormal": extracted["is_abnormal"],
    }, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    # 4-b. PPTX — 바이너리 대신 슬라이드 텍스트(POI 비교 대상).
    from pptx import Presentation
    pptx_path = workp / "golden.pptx"
    sections = [
        {"title": "제안 개요", "content": "골든 캡처 고정 본문 — 값 임의 변경 금지",
         "sources": ["spec/01", "plan FN-1"]},
        {"title": "세부 계획", "content": "두 번째 섹션 고정 본문", "sources": []},
    ]
    build_pptx(sections, scoring["criteria"], str(pptx_path), institution_name="노원구")
    slides = []
    for slide in Presentation(str(pptx_path)).slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.extend(p.text for p in shape.text_frame.paragraphs if p.text)
        slides.append(texts)
    (art_dir / "pptx_slides.json").write_text(
        json.dumps(slides, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    manifest = {
        "captured_endpoints": captured,
        "seeded_institutions": seeded,
        "fts_index": {"files": stats.get("files"), "chunks": stats.get("chunks")},
        "note": "재캡처는 py -3.14 -m golden.capture — git diff가 변화 목록이다",
    }
    (GOLDEN / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    shutil.rmtree(work, ignore_errors=True)
    print(f"golden: {len(captured)}개 API + 산출물 2종 캡처 완료 (기관 {seeded}건 시드)")


if __name__ == "__main__":
    main()
