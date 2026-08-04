"""확장자별 코퍼스 파서 레지스트리.

v1은 .txt(UTF-8)만 지원한다. rfp(PDF)·reports(HTML/DOCX)는 여기에 확장자 항목을
추가하는 것으로 확장한다 — 스펙 §④·§⑦ 참조.
"""

from __future__ import annotations

import json
from pathlib import Path


def _parse_txt(path: Path) -> str | None:
    """UTF-8로 읽되, 실패는 조용히 None — 인코딩 검증은 corpus_validator의 몫이다."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return None


def _parse_pptx(path: Path) -> str | None:
    """제안서 슬라이드의 글자를 뽑는다 — 도형 텍스트 + 표 셀.

    계획 F Task 5에서 추가했다. 아카이브에 남는 진짜 산출물이 제안서 `.pptx`인데,
    파서가 없으면 "완료 후 산출물이 지식 탭에서 검색된다"(스펙 §② 17)가 `rfp_text.txt`
    한 건짜리 약속이 돼버린다. python-pptx는 이미 의존성에 있다(pptx_builder).

    깨진 파일 하나로 색인 전체가 멈추면 안 되므로 실패는 조용히 None이다 —
    `.txt`의 인코딩 실패와 같은 취급.
    """
    try:
        from pptx import Presentation
    except ImportError:      # pragma: no cover - 의존성이 빠진 환경
        return None

    try:
        presentation = Presentation(str(path))
    except Exception:
        return None

    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    # 문단 경계를 빈 줄로 — chunker가 빈 줄 기준으로 나눈다.
    return "\n\n".join(parts) if parts else None


def _parse_json(path: Path) -> str | None:
    """구조화 산출물(`rfp_scoring.json`·`coverage_map.json`)을 검색 가능한 줄로 편다.

    배점표 항목명("금고 운영 실적", "지역사회 기여" 같은 평가 기준)은 실제로 가장
    자주 찾는 것 중 하나인데, 파서가 없어 **허용목록에 이름만 올라 있고 영영 색인될 수
    없는 상태**였다(`ARCHIVE_INDEXABLE_NAMES`).

    원문 JSON을 그대로 넣지 않는 이유: `{"criteria": [{"name":` 같은 문법 부스러기가
    본문에 섞이면 스니펫이 읽히지 않고 매치도 지저분해진다. 키는 **라벨**로 붙이고
    값만 남긴다.

    아카이브 밖에서는 사실상 무해하다 — `corpus/`에 `.json`이 하나도 없고, 아카이브
    안에서는 허용목록이 `tasks_dump.json`(대화 원문)을 막는다.
    """
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError, OSError):
        return None

    lines: list[str] = []

    def walk(node, label: str) -> None:
        if isinstance(node, dict):
            for key, value in node.items():
                walk(value, f"{label} {key}".strip())
        elif isinstance(node, list):
            for item in node:
                walk(item, label)
        elif node is not None and node != "":
            lines.append(f"{label}: {node}" if label else str(node))

    walk(data, "")
    return "\n\n".join(lines) if lines else None


PARSERS = {
    ".txt": _parse_txt,
    ".pptx": _parse_pptx,
    ".json": _parse_json,
}


def parse_file(path: Path) -> str | None:
    """지원 확장자면 본문 텍스트를, 아니면(또는 디코딩 실패면) None을 돌려준다."""
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        return None
    return parser(path)
