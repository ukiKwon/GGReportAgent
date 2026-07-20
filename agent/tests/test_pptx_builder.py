import os

from pptx import Presentation

from agent.nodes.pptx_builder import build_pptx


def test_build_pptx_creates_title_and_section_slides(tmp_path):
    sections = [
        {"scoring_item": "신용도", "title": "1. 신용도 및 재무구조", "content": "본문 내용 1", "sources": ["spec/01"]},
        {"scoring_item": "예금금리", "title": "2. 예금 및 대출금리", "content": "본문 내용 2", "sources": ["spec/02"]},
    ]
    scoring_table = [
        {"category": "신용도", "item": "외부평가", "score": 8, "description": None},
        {"category": "예금금리", "item": "정기예금금리", "score": 21, "description": None},
    ]
    out_path = str(tmp_path / "test.pptx")

    build_pptx(sections, scoring_table, out_path)

    assert os.path.isfile(out_path)
    prs = Presentation(out_path)
    # title slide + scoring table slide + 2 section slides = 4 minimum
    assert len(prs.slides) >= 4
    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "신용도 및 재무구조" in all_text
    assert "본문 내용 1" in all_text


def test_build_pptx_splices_archive_slides_as_separate_section(tmp_path):
    # Build a fake archive PPTX with one distinguishable slide
    archive_path = str(tmp_path / "archive.pptx")
    archive_prs = Presentation()
    slide = archive_prs.slides.add_slide(archive_prs.slide_layouts[1])
    slide.shapes.title.text = "과거 제안 슬라이드 마커"
    archive_prs.save(archive_path)

    sections = [{"scoring_item": "신용도", "title": "1. 신용도", "content": "내용", "sources": ["spec/01"]}]
    scoring_table = [{"category": "신용도", "item": "평가", "score": 10, "description": None}]
    out_path = str(tmp_path / "test.pptx")

    build_pptx(sections, scoring_table, out_path, archive_pptx_path=archive_path)

    prs = Presentation(out_path)
    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "과거 유사제안 참고" in all_text  # section divider slide
    assert "과거 제안 슬라이드 마커" in all_text  # copied original content


def test_build_pptx_without_archive_has_no_reference_section(tmp_path):
    sections = [{"scoring_item": "신용도", "title": "1. 신용도", "content": "내용", "sources": ["spec/01"]}]
    scoring_table = [{"category": "신용도", "item": "평가", "score": 10, "description": None}]
    out_path = str(tmp_path / "test.pptx")

    build_pptx(sections, scoring_table, out_path, archive_pptx_path=None)

    prs = Presentation(out_path)
    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "과거 유사제안 참고" not in all_text
