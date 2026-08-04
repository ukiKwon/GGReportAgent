from agent.retrieval.parsers import parse_file


def test_parse_txt_reads_utf8(tmp_path):
    path = tmp_path / "00_인덱스.txt"
    path.write_text("도봉구 사업목록", encoding="utf-8")
    assert parse_file(path) == "도봉구 사업목록"


def test_parse_txt_returns_none_for_non_utf8(tmp_path):
    path = tmp_path / "broken.txt"
    path.write_bytes("한글".encode("cp949"))
    assert parse_file(path) is None


def test_unsupported_extension_returns_none(tmp_path):
    path = tmp_path / "공고문.pdf"
    path.write_bytes(b"%PDF-1.4")
    assert parse_file(path) is None


def test_extension_match_is_case_insensitive(tmp_path):
    path = tmp_path / "NOTE.TXT"
    path.write_text("내용", encoding="utf-8")
    assert parse_file(path) == "내용"


# ── .pptx (계획 F Task 5) ────────────────────────────────────────────────
# 아카이브에 남는 진짜 산출물이 제안서 pptx다. 파서가 없으면 "완료 후 산출물이
# 지식 탭에서 검색된다"(스펙 §② 17)가 rfp_text.txt 한 건짜리 약속이 돼버린다.


def _write_pptx(path, title, body, table_rows=None):
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    if table_rows:
        shape = slide.shapes.add_table(
            len(table_rows), len(table_rows[0]),
            Inches(1), Inches(4), Inches(6), Inches(1),
        )
        for r, row in enumerate(table_rows):
            for c, value in enumerate(row):
                shape.table.cell(r, c).text = value
    presentation.save(str(path))


def test_pptx에서_도형_텍스트를_뽑는다(tmp_path):
    path = tmp_path / "제안서.pptx"
    _write_pptx(path, "도봉구 금고 제안", "청년 창업 연계 금융상품")

    text = parse_file(path)

    assert "도봉구 금고 제안" in text
    assert "청년 창업 연계 금융상품" in text


def test_pptx의_표_셀도_뽑는다(tmp_path):
    """배점표가 표로 들어가는 일이 많아, 표를 빼면 정작 쓸모 있는 게 안 잡힌다."""
    path = tmp_path / "배점.pptx"
    _write_pptx(path, "평가 배점", "요약", [["항목", "배점"], ["금고 운영 실적", "20"]])

    text = parse_file(path)

    assert "금고 운영 실적 | 20" in text


def test_깨진_pptx는_조용히_건너뛴다(tmp_path):
    """파일 하나 때문에 색인 전체가 멈추면 안 된다 — txt 인코딩 실패와 같은 취급."""
    path = tmp_path / "깨짐.pptx"
    path.write_bytes(b"not a real pptx")
    assert parse_file(path) is None


def test_빈_pptx는_None이다(tmp_path):
    from pptx import Presentation

    path = tmp_path / "빈.pptx"
    Presentation().save(str(path))
    assert parse_file(path) is None
