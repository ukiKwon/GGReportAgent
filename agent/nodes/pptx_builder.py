import copy
import os

from pptx import Presentation


def _add_title_slide(prs, institution_name):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{institution_name} 제안서"


def _add_scoring_table_slide(prs, scoring_table):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "평가 배점표"
    body = slide.placeholders[1].text_frame
    body.text = f"{scoring_table[0]['category']}: {scoring_table[0]['item']} ({scoring_table[0]['score']}점)"
    for item in scoring_table[1:]:
        p = body.add_paragraph()
        p.text = f"{item['category']}: {item['item']} ({item['score']}점)"


def _add_section_slide(prs, section):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = section["title"]
    body = slide.placeholders[1].text_frame
    body.text = section["content"]
    if section.get("sources"):
        p = body.add_paragraph()
        p.text = f"근거자료: {', '.join(section['sources'])}"
        p.level = 1


def _add_archive_reference_section(prs, archive_pptx_path):
    divider = prs.slides.add_slide(prs.slide_layouts[2])
    divider.shapes.title.text = "과거 유사제안 참고"

    archive_prs = Presentation(archive_pptx_path)
    blank_layout = prs.slide_layouts[6]
    for archive_slide in archive_prs.slides:
        new_slide = prs.slides.add_slide(blank_layout)
        for shape in archive_slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape._element))


def build_pptx(sections, scoring_table, output_path, archive_pptx_path=None, institution_name="기관"):
    prs = Presentation()
    _add_title_slide(prs, institution_name)
    _add_scoring_table_slide(prs, scoring_table)
    for section in sections:
        _add_section_slide(prs, section)
    if archive_pptx_path:
        _add_archive_reference_section(prs, archive_pptx_path)
    prs.save(output_path)


def pptx_builder_node(state: dict) -> dict:
    institution_name = state.get("institution_name", "기관")
    output_path = f"report_new/{institution_name}/{institution_name}_제안서.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    build_pptx(
        state["sections"],
        state["scoring_table"],
        output_path,
        archive_pptx_path=state.get("archive_pptx_path"),
        institution_name=institution_name,
    )
    return {"pptx_path": output_path}
