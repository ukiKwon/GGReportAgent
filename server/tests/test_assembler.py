import pytest
from pptx import Presentation

from server.assembler import assemble_deliverable
from server.db import init_db
from server.teams import AUTHORING_TEAMS


@pytest.fixture
def conn(tmp_path):
    connection = init_db(str(tmp_path / "test.db"))
    connection.execute(
        "INSERT INTO institutions (institution_id, name_ko, stage) VALUES ('mapo', '마포구', 6)"
    )
    connection.execute(
        "INSERT INTO bid_cases (bid_case_id, institution_id) VALUES ('bc-1', 'mapo')"
    )
    # ⚠️ 팀 이름을 손으로 적지 않는다. 예전에는 여기가 `"IT"`(옛 이름)였고 구현도 같은
    # 옛 이름을 쓰고 있어서, **둘이 나란히 틀린 채로 통과**했다 — 실제 운영에서는
    # `전산` 팀 초안이 취합에서 통째로 빠지고 있었는데 이 테스트가 못 잡았다.
    for team in AUTHORING_TEAMS:
        draft = f"{team} 초안 본문"
        connection.execute(
            """INSERT INTO tasks (task_id, bid_case_id, team, status, draft_content)
               VALUES (?, 'bc-1', ?, '2차완료', ?)""",
            (f"task-{team}", team, draft),
        )
    connection.commit()
    yield connection
    connection.close()


def _slide_texts(path):
    return [
        "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        for slide in Presentation(path).slides
    ]


def test_assemble_writes_one_slide_per_authoring_team_in_order(conn, tmp_path):
    """작성 3팀이 **하나도 빠짐없이**, 정해진 순서로 들어간다.

    `AUTHORING_TEAMS`를 그대로 돌려 비교하므로, 구현이 팀 이름을 따로 박아 두면
    (예전의 `TEAM_ORDER = [..., "IT", ...]`) 그 자리에서 실패한다.
    """
    path = assemble_deliverable(conn, "bc-1", output_root=str(tmp_path / "out"))

    texts = _slide_texts(path)
    team_slides = [t for t in texts if "초안 본문" in t]
    assert len(team_slides) == len(AUTHORING_TEAMS)
    for slide, team in zip(team_slides, AUTHORING_TEAMS):
        assert f"{team} 초안 본문" in slide


def test_assemble_does_not_silently_drop_a_team(conn, tmp_path):
    """빠진 팀이 있으면 **오류 없이 슬라이드 한 장이 없어질 뿐**이라 따로 못 박는다.

    실제로 그렇게 `전산` 팀 초안이 계속 빠지고 있었다(2026-08-27 Java 이관 중 발견).
    위 순서 테스트와 달리 이건 "무엇이 빠졌는지"를 실패 메시지에 그대로 보여준다.
    """
    path = assemble_deliverable(conn, "bc-1", output_root=str(tmp_path / "out"))

    deck = "\n".join(_slide_texts(path))
    missing = [team for team in AUTHORING_TEAMS if f"{team} 파트" not in deck]
    assert not missing, f"취합에서 빠진 팀: {missing}"


def test_assemble_titles_the_deck_with_the_institution_name(conn, tmp_path):
    path = assemble_deliverable(conn, "bc-1", output_root=str(tmp_path / "out"))

    assert "마포구" in _slide_texts(path)[0]


def test_assemble_records_pptx_path_on_the_institution(conn, tmp_path):
    path = assemble_deliverable(conn, "bc-1", output_root=str(tmp_path / "out"))

    stored = conn.execute(
        "SELECT pptx_path FROM institutions WHERE institution_id = 'mapo'"
    ).fetchone()["pptx_path"]
    assert stored == path


def test_assemble_raises_when_bid_case_is_missing(conn, tmp_path):
    with pytest.raises(KeyError):
        assemble_deliverable(conn, "bc-nope", output_root=str(tmp_path / "out"))
